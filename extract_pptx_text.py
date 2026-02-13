# -*- coding: utf-8 -*-
"""Extract all text from a .pptx to find company names."""
import sys
import os

def main():
    path = r"c:\Users\ryanx\OneDrive\桌面\A股上市公司2021年-2025年同行业对比-刀刀.pptx"
    if len(sys.argv) > 1:
        path = sys.argv[1]
    if not os.path.isfile(path):
        print("File not found:", path, file=sys.stderr)
        sys.exit(1)
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        print("Run: pip install python-pptx", file=sys.stderr)
        sys.exit(2)
    prs = Presentation(path)
    all_text = []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                all_text.append(shape.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text:
                            all_text.append(cell.text.strip())
    text = "\n".join(all_text)
    # Ensure stdout uses utf-8 for Windows console
    import io
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
    print(text)

if __name__ == "__main__":
    main()
